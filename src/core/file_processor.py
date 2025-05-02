from pathlib import Path
import logging
from typing import Optional, Dict, Any
from docx import Document
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
import json
import csv
import chardet

logger = logging.getLogger(__name__)

class FileProcessor:
    """Handles processing of various file types with content extraction"""
    
    SUPPORTED_EXTENSIONS = {
        '.txt': 'text',
        '.docx': 'word',
        '.pdf': 'pdf',
        '.pptx': 'powerpoint',
        '.xlsx': 'excel',
        '.csv': 'csv',
        '.json': 'json',
        '.md': 'text'
    }

    @classmethod
    def can_process(cls, file_path: str) -> bool:
        """Check if the file type is supported"""
        return Path(file_path).suffix.lower() in cls.SUPPORTED_EXTENSIONS

    def extract_text(self, file_path: Path) -> Optional[str]:
        """Extract text content from supported file types"""
        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            return None

        try:
            file_type = self.SUPPORTED_EXTENSIONS.get(file_path.suffix.lower())
            if not file_type:
                logger.warning(f"Unsupported file type: {file_path.suffix}")
                return None

            method_name = f"_process_{file_type}"
            if hasattr(self, method_name):
                return getattr(self, method_name)(file_path)
            
            return None

        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            return None

    def _process_text(self, file_path: Path) -> str:
        """Process text files with encoding detection"""
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            detected = chardet.detect(raw_data)
            encoding = detected['encoding'] or 'utf-8'
            
        return file_path.read_text(encoding=encoding)

    def _process_word(self, file_path: Path) -> str:
        """Process Word documents"""
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = []
        
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                if any(row_data):  # Only include non-empty rows
                    table_data.append(row_data)
            if table_data:
                tables.append(table_data)

        content = "\n\n".join(paragraphs)
        if tables:
            content += "\n\nTables:\n"
            for i, table in enumerate(tables, 1):
                content += f"\nTable {i}:\n"
                content += "\n".join([" | ".join(row) for row in table])
                content += "\n"

        return content

    def _process_pdf(self, file_path: Path) -> str:
        """Process PDF files"""
        reader = PdfReader(file_path)
        text_content = []
        
        for page in reader.pages:
            text = page.extract_text()
            if text.strip():
                text_content.append(text)
        
        return "\n\n".join(text_content)

    def _process_powerpoint(self, file_path: Path) -> str:
        """Process PowerPoint presentations"""
        prs = Presentation(file_path)
        content = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_content = [f"Slide {i}:"]
            
            if slide.shapes.title:
                slide_content.append(f"Title: {slide.shapes.title.text}")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:  # Skip title as it's already added
                        slide_content.append(shape.text)
            
            content.append("\n".join(slide_content))
        
        return "\n\n".join(content)

    def _process_excel(self, file_path: Path) -> str:
        """Process Excel files"""
        dfs = pd.read_excel(file_path, sheet_name=None)
        content = []
        
        for sheet_name, df in dfs.items():
            content.append(f"Sheet: {sheet_name}")
            content.append(df.to_string(index=False))
        
        return "\n\n".join(content)

    def _process_csv(self, file_path: Path) -> str:
        """Process CSV files with encoding detection"""
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            detected = chardet.detect(raw_data)
            encoding = detected['encoding'] or 'utf-8'
        
        df = pd.read_csv(file_path, encoding=encoding)
        return df.to_string(index=False)

    def _process_json(self, file_path: Path) -> str:
        """Process JSON files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return json.dumps(data, indent=2)

    def get_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from the file"""
        stats = file_path.stat()
        return {
            "name": file_path.name,
            "size": stats.st_size,
            "created": stats.st_ctime,
            "modified": stats.st_mtime,
            "type": self.SUPPORTED_EXTENSIONS.get(file_path.suffix.lower(), "unknown"),
            "extension": file_path.suffix.lower()
        } 