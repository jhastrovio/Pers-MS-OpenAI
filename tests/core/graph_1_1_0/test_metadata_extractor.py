"""
Test for the metadata extractor module.
"""

import os
import pytest
import io
from core.graph_1_1_0.metadata_extractor import MetadataExtractor
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from pypdf import PdfWriter

class TestMetadataExtractor:
    """Test suite for the MetadataExtractor class."""

    def test_extract_word_metadata(self):
        """Test extracting metadata from Word documents."""
        # Create a simple Word document
        doc = Document()
        doc.add_heading('Test Document', 0)
        doc.add_paragraph('This is a test paragraph for metadata extraction.')
        doc.add_paragraph('Another test paragraph.')
        
        # Save to bytes
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        content = file_stream.getvalue()
        
        # Test extraction
        metadata = MetadataExtractor.extract_metadata(
            content, 
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
        
        # Verify metadata fields
        assert 'paragraph_count' in metadata
        assert metadata['paragraph_count'] >= 2
        assert 'text_content' in metadata
        assert 'Test Document' in metadata['text_content']
        assert 'This is a test paragraph' in metadata['text_content']

    def test_extract_ppt_metadata(self):
        """Test extracting metadata from PowerPoint files."""
        # Create a simple PowerPoint presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Presentation"
        
        # Add a second slide with content
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Second Slide"
        content = slide.placeholders[1]
        content.text = "This is test content for extraction."
        
        # Save to bytes
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        content = file_stream.getvalue()
        
        # Test extraction
        metadata = MetadataExtractor.extract_metadata(
            content, 
            'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
        # Verify metadata fields
        assert 'slide_count' in metadata
        assert metadata['slide_count'] == 2
        assert 'slide_titles' in metadata
        assert len(metadata['slide_titles']) == 2
        assert 'text_content' in metadata
        assert 'Test Presentation' in metadata['text_content']
        assert 'Second Slide' in metadata['text_content']

    def test_extract_excel_metadata(self):
        """Test extracting metadata from Excel files."""
        # Create a simple Excel workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Test Sheet"
        
        # Add header row
        ws['A1'] = 'Name'
        ws['B1'] = 'Age'
        ws['C1'] = 'Location'
        
        # Add data rows
        ws['A2'] = 'John'
        ws['B2'] = 30
        ws['C2'] = 'New York'
        
        ws['A3'] = 'Jane'
        ws['B3'] = 25
        ws['C3'] = 'London'
        
        # Save to bytes
        file_stream = io.BytesIO()
        wb.save(file_stream)
        file_stream.seek(0)
        content = file_stream.getvalue()
        
        # Test extraction
        metadata = MetadataExtractor.extract_metadata(
            content, 
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )
        
        # Verify metadata fields
        assert 'sheet_count' in metadata
        assert metadata['sheet_count'] == 1
        assert 'sheet_names' in metadata
        assert 'Test Sheet' in metadata['sheet_names']
        assert 'column_names' in metadata
        assert len(metadata['column_names']) == 3
        assert 'text_content' in metadata
        assert 'Name' in metadata['text_content']
        assert 'John' in metadata['text_content']

    def test_extract_text_metadata(self):
        """Test extracting metadata from text files."""
        text_content = "This is a test text file.\nIt has multiple lines.\nThird line for testing."
        content = text_content.encode('utf-8')
        
        # Test extraction
        metadata = MetadataExtractor.extract_metadata(content, 'text/plain')
        
        # Verify metadata fields
        assert 'line_count' in metadata
        assert metadata['line_count'] == 3
        assert 'word_count' in metadata
        assert metadata['word_count'] >= 12
        assert 'text_content' in metadata
        assert 'This is a test text file.' in metadata['text_content']

    def test_extract_html_metadata(self):
        """Test extracting metadata from HTML files."""
        html_content = """
        <!DOCTYPE html>
        <html>
        <head>
            <title>Test HTML Document</title>
            <meta name="author" content="Test Author">
            <meta name="description" content="This is a test HTML document">
        </head>
        <body>
            <h1>Test Heading</h1>
            <p>This is a test paragraph.</p>
            <p>Another test paragraph with <b>formatted</b> text.</p>
        </body>
        </html>
        """
        content = html_content.encode('utf-8')
        
        # Test extraction
        metadata = MetadataExtractor.extract_metadata(content, 'text/html')
        
        # Verify metadata fields
        assert 'title' in metadata
        assert metadata['title'] == 'Test HTML Document'
        assert 'author' in metadata
        assert metadata['author'] == 'Test Author'
        assert 'description' in metadata
        assert metadata['description'] == 'This is a test HTML document'
        assert 'text_content' in metadata
        assert 'Test Heading' in metadata['text_content']
        assert 'This is a test paragraph' in metadata['text_content']

    def test_unsupported_content_type(self):
        """Test behavior with unsupported content type."""
        content = b'Test content'
        
        # Test extraction with unsupported type
        metadata = MetadataExtractor.extract_metadata(content, 'application/unknown')
        
        # Should return empty dict
        assert metadata == {} 