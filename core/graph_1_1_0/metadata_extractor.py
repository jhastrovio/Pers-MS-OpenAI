"""
Metadata extraction utilities for Graph API documents.

This module provides functionality to extract metadata from different file types
retrieved from Microsoft Graph API.
"""

import io
from typing import Dict, Any
import pandas as pd
from pptx import Presentation
from openpyxl import load_workbook
import docx
from pypdf import PdfReader
import chardet
from bs4 import BeautifulSoup
import os
from datetime import datetime

class MetadataExtractor:
    """Handles metadata extraction from different file types retrieved from Graph API."""
    
    @staticmethod
    def extract_metadata(content: bytes, content_type: str) -> Dict[str, Any]:
        """Extract metadata from content based on its type.
        
        Args:
            content: Raw content bytes
            content_type: MIME type of the content
            
        Returns:
            Dictionary containing extracted metadata
        """
        content_type = content_type.lower()
        
        try:
            if content_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
                return MetadataExtractor._extract_ppt_metadata(content)
            elif content_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel', 'text/csv']:
                return MetadataExtractor._extract_excel_metadata(content, content_type)
            elif content_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
                return MetadataExtractor._extract_word_metadata(content)
            elif content_type == 'application/pdf':
                return MetadataExtractor._extract_pdf_metadata(content)
            elif content_type == 'text/plain':
                return MetadataExtractor._extract_text_metadata(content)
            elif content_type == 'text/html':
                return MetadataExtractor._extract_html_metadata(content)
            else:
                return {}
        except Exception as e:
            return {}
    
    @staticmethod
    def _extract_ppt_metadata(content: bytes) -> Dict[str, Any]:
        """Extract metadata from PowerPoint content."""
        try:
            ppt_file = io.BytesIO(content)
            prs = Presentation(ppt_file)
            
            # Convert modified datetime to string if it exists
            last_modified = ''
            if hasattr(prs.core_properties, 'modified') and prs.core_properties.modified:
                if isinstance(prs.core_properties.modified, datetime):
                    last_modified = prs.core_properties.modified.isoformat()
                else:
                    last_modified = str(prs.core_properties.modified)
            
            metadata = {
                'title': prs.core_properties.title or '',
                'author': prs.core_properties.author or '',
                'last_modified': last_modified,
                'slide_count': len(prs.slides),
                'slide_titles': []
            }
            
            # Extract slide titles
            for slide in prs.slides:
                title = ''
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        title = shape.text.strip()
                        break
                metadata['slide_titles'].append(title)
            
            # Extract text content
            text_content = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                text_content.append(" ".join(slide_text))
            
            metadata['text_content'] = "\n\n".join(text_content)
            
            return metadata
        except Exception as e:
            return {
                'title': '',
                'author': '',
                'last_modified': '',
                'slide_count': 0,
                'slide_titles': [],
                'text_content': ''
            }
    
    @staticmethod
    def _extract_excel_metadata(content: bytes, content_type: str) -> Dict[str, Any]:
        """Extract metadata from Excel/CSV content."""
        try:
            if content_type == 'text/csv':
                # Handle CSV files
                csv_file = io.StringIO(content.decode('utf-8', errors='replace'))
                df = pd.read_csv(csv_file)
                metadata = {
                    'sheet_count': 1,
                    'sheet_names': ['Sheet1'],
                    'row_count': len(df),
                    'column_count': len(df.columns),
                    'column_names': df.columns.tolist()
                }
                
                # Get a sample of the content as text
                metadata['text_content'] = df.head(20).to_string(index=False)
            else:
                # Handle Excel files
                excel_file = io.BytesIO(content)
                wb = load_workbook(excel_file, read_only=True, data_only=True)
                metadata = {
                    'sheet_count': len(wb.sheetnames),
                    'sheet_names': wb.sheetnames,
                    'row_count': 0,
                    'column_count': 0,
                    'column_names': [],
                    'text_content': ''
                }
                
                # Get dimensions of the first sheet
                if wb.sheetnames:
                    ws = wb[wb.sheetnames[0]]
                    metadata['row_count'] = ws.max_row
                    metadata['column_count'] = ws.max_column
                    
                    # Extract headers and a sample of content
                    headers = []
                    sample_rows = []
                    row_count = 0
                    
                    for row in ws.iter_rows(max_row=20):  # Get first 20 rows for sample
                        row_values = [str(cell.value) if cell.value is not None else '' for cell in row]
                        if row_count == 0:
                            headers = row_values
                            metadata['column_names'] = headers
                        else:
                            sample_rows.append(row_values)
                        row_count += 1
                    
                    # Format as string for text content
                    content_lines = []
                    if headers:
                        content_lines.append(" | ".join(headers))
                        content_lines.append("-" * (sum(len(h) for h in headers) + 3 * len(headers)))
                    
                    for row in sample_rows:
                        content_lines.append(" | ".join(row))
                    
                    metadata['text_content'] = "\n".join(content_lines)
            
            return metadata
        except Exception as e:
            return {
                'sheet_count': 0,
                'sheet_names': [],
                'row_count': 0,
                'column_count': 0,
                'column_names': [],
                'text_content': ''
            }
    
    @staticmethod
    def _extract_word_metadata(content: bytes) -> Dict[str, Any]:
        """Extract metadata from Word document."""
        try:
            word_file = io.BytesIO(content)
            doc = docx.Document(word_file)
            
            # Get core properties
            core_props = doc.core_properties
            
            # Convert modified datetime to string if it exists
            last_modified = ''
            if hasattr(core_props, 'modified') and core_props.modified:
                if isinstance(core_props.modified, datetime):
                    last_modified = core_props.modified.isoformat()
                else:
                    last_modified = str(core_props.modified)
            
            metadata = {
                'title': core_props.title if hasattr(core_props, 'title') else '',
                'author': core_props.author if hasattr(core_props, 'author') else '',
                'last_modified': last_modified,
                'page_count': len(doc.sections),
                'paragraph_count': len(doc.paragraphs),
                'word_count': sum(len(paragraph.text.split()) for paragraph in doc.paragraphs if paragraph.text),
                'text_content': '\n\n'.join([paragraph.text for paragraph in doc.paragraphs if paragraph.text])
            }
            
            return metadata
        except Exception as e:
            return {
                'title': '',
                'author': '',
                'last_modified': '',
                'page_count': 0,
                'paragraph_count': 0,
                'word_count': 0,
                'text_content': ''
            }
    
    @staticmethod
    def _extract_pdf_metadata(content: bytes) -> Dict[str, Any]:
        """Extract metadata from PDF document."""
        try:
            pdf_file = io.BytesIO(content)
            pdf = PdfReader(pdf_file)
            
            # Extract info from document information dictionary
            info = pdf.metadata
            
            # Convert any datetime objects to strings
            modification_date = None
            creation_date = None
            
            if info and hasattr(info, 'modification_date') and info.modification_date:
                if isinstance(info.modification_date, datetime):
                    modification_date = info.modification_date.isoformat()
                else:
                    modification_date = str(info.modification_date)
                    
            if info and hasattr(info, 'creation_date') and info.creation_date:
                if isinstance(info.creation_date, datetime):
                    creation_date = info.creation_date.isoformat()
                else:
                    creation_date = str(info.creation_date)
            
            metadata = {
                'title': info.title if info and hasattr(info, 'title') else '',
                'author': info.author if info and hasattr(info, 'author') else '',
                'last_modified': modification_date,
                'creation_date': creation_date,
                'producer': info.producer if info and hasattr(info, 'producer') else '',
                'page_count': len(pdf.pages),
                'text_content': ''
            }
            
            # Extract text from first few pages
            text_content = []
            max_pages = min(5, len(pdf.pages))  # Limit to first 5 pages
            for i in range(max_pages):
                page_text = pdf.pages[i].extract_text()
                if page_text:
                    text_content.append(page_text)
            
            metadata['text_content'] = '\n\n'.join(text_content)
            
            return metadata
        except Exception as e:
            return {
                'title': '',
                'author': '',
                'last_modified': '',
                'creation_date': '',
                'producer': '',
                'page_count': 0,
                'text_content': ''
            }
    
    @staticmethod
    def _extract_text_metadata(content: bytes) -> Dict[str, Any]:
        """Extract metadata from plain text file."""
        try:
            # Detect encoding
            detection = chardet.detect(content)
            encoding = detection['encoding'] if detection['encoding'] else 'utf-8'
            
            # Decode content
            text = content.decode(encoding, errors='replace')
            
            # Count lines, words, and characters
            lines = text.split('\n')
            words = text.split()
            
            metadata = {
                'encoding': encoding,
                'confidence': detection['confidence'],
                'line_count': len(lines),
                'word_count': len(words),
                'char_count': len(text),
                'text_content': text[:5000]  # Limit text content preview
            }
            
            return metadata
        except Exception as e:
            return {
                'encoding': 'unknown',
                'confidence': 0,
                'line_count': 0,
                'word_count': 0,
                'char_count': 0,
                'text_content': ''
            }
    
    @staticmethod
    def _extract_html_metadata(content: bytes) -> Dict[str, Any]:
        """Extract metadata from HTML file."""
        try:
            # Detect encoding
            detection = chardet.detect(content)
            encoding = detection['encoding'] if detection['encoding'] else 'utf-8'
            
            # Decode content
            html_text = content.decode(encoding, errors='replace')
            
            # Parse HTML
            soup = BeautifulSoup(html_text, 'html.parser')
            
            # Extract metadata
            title = soup.title.string if soup.title else ''
            
            # Get meta tags
            meta_tags = {}
            for tag in soup.find_all('meta'):
                if tag.get('name') and tag.get('content'):
                    meta_tags[tag['name']] = tag['content']
            
            # Extract visible text
            [s.extract() for s in soup(['style', 'script', 'head', '[document]', 'title'])]
            visible_text = ' '.join(soup.stripped_strings)
            
            metadata = {
                'title': title,
                'author': meta_tags.get('author', ''),
                'description': meta_tags.get('description', ''),
                'keywords': meta_tags.get('keywords', ''),
                'encoding': encoding,
                'confidence': detection['confidence'],
                'text_content': visible_text[:5000]  # Limit text content preview
            }
            
            return metadata
        except Exception as e:
            return {
                'title': '',
                'author': '',
                'description': '',
                'keywords': '',
                'encoding': 'unknown',
                'confidence': 0,
                'text_content': ''
            } 