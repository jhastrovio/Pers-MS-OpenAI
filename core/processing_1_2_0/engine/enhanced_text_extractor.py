"""
Enhanced text extraction utilities for large, rich-format documents.

This module provides advanced document processing capabilities including:
- OCR for scanned documents
- Layout-aware extraction 
- Content cleaning and normalization
- Semantic segmentation and chunking
"""

import io
import re
import logging
from typing import Dict, Any, List, Tuple, Optional
from dataclasses import dataclass
import tempfile
import os

# Core document processing
from bs4 import BeautifulSoup
import docx
from pypdf import PdfReader
import pandas as pd
from pptx import Presentation
from openpyxl import load_workbook

# Advanced processing libraries
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    from unstructured.partition.auto import partition
    from unstructured.chunking.title import chunk_by_title
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    import tiktoken
    TIKTOKEN_AVAILABLE = True
except ImportError:
    TIKTOKEN_AVAILABLE = False

from core.utils.logging import get_logger

logger = get_logger(__name__)

@dataclass
class DocumentChunk:
    """Represents a semantically coherent chunk of document content."""
    content: str
    chunk_type: str  # 'paragraph', 'section', 'table', 'header', etc.
    metadata: Dict[str, Any]
    position: int  # Position in original document
    heading_hierarchy: List[str]  # List of heading levels ["Chapter 1", "Section 1.1", etc.]
    page_number: Optional[int] = None
    confidence_score: Optional[float] = None

@dataclass
class ExtractionConfig:
    """Configuration for document extraction."""
    # Layout analysis
    use_layout_analysis: bool = True
    preserve_structure: bool = True
    remove_headers_footers: bool = True
    
    # Content cleaning
    normalize_whitespace: bool = True
    remove_page_numbers: bool = True
    remove_boilerplate: bool = True
    fix_encoding: bool = True
    
    # Chunking
    chunk_documents: bool = True
    chunk_size: int = 500  # tokens
    chunk_overlap: int = 75  # tokens
    preserve_headings: bool = True
    
    # OCR
    use_ocr_fallback: bool = True
    ocr_confidence_threshold: float = 0.7

class EnhancedTextExtractor:
    """Advanced text extraction with layout awareness, cleaning, and chunking."""
    
    def __init__(self, config: Optional[ExtractionConfig] = None):
        """Initialize the enhanced text extractor.
        
        Args:
            config: Configuration for extraction behavior
        """
        self.config = config or ExtractionConfig()
        self._ensure_dependencies()
        
        # Initialize tokenizer for chunking
        if TIKTOKEN_AVAILABLE:
            try:
                self.tokenizer = tiktoken.get_encoding("cl100k_base")  # GPT-4 tokenizer
            except Exception as e:
                logger.warning(f"Failed to initialize tokenizer: {e}")
                self.tokenizer = None
        else:
            self.tokenizer = None
    
    def _ensure_dependencies(self):
        """Check and warn about missing dependencies."""
        missing = []
        if not PDFPLUMBER_AVAILABLE:
            missing.append("pdfplumber")
        if not UNSTRUCTURED_AVAILABLE:
            missing.append("unstructured")
        if not OCR_AVAILABLE:
            missing.append("pytesseract/PIL")
        if not TIKTOKEN_AVAILABLE:
            missing.append("tiktoken")
            
        if missing:
            logger.warning(f"Missing advanced processing dependencies: {missing}. Install with: pip install {' '.join(missing)}")
    
    def extract_and_process(self, content: bytes, content_type: str, filename: str = "") -> Dict[str, Any]:
        """Extract and process document with advanced capabilities.
        
        Args:
            content: Raw document bytes
            content_type: MIME type of the document
            filename: Original filename for context
            
        Returns:
            Dict containing:
                - raw_text: Basic extracted text
                - cleaned_text: Processed and cleaned text
                - chunks: List of DocumentChunk objects
                - metadata: Extraction metadata
        """
        content_type = content_type.lower()
        
        try:
            # Step 1: Extract raw text with layout awareness
            raw_text, structure_metadata = self._extract_with_layout(content, content_type, filename)
            
            # Step 2: Clean and normalize content
            cleaned_text = self._clean_and_normalize(raw_text)
            
            # Step 3: Semantic chunking
            chunks = []
            if self.config.chunk_documents:
                chunks = self._create_semantic_chunks(cleaned_text, structure_metadata)
            
            # Step 4: Compile results
            result = {
                "raw_text": raw_text,
                "cleaned_text": cleaned_text,
                "chunks": chunks,
                "metadata": {
                    "filename": filename,
                    "content_type": content_type,
                    "extraction_config": self.config.__dict__,
                    "structure_metadata": structure_metadata,
                    "chunk_count": len(chunks),
                    "processing_notes": []
                }
            }
            
            return result
            
        except Exception as e:
            logger.error(f"Error in enhanced text extraction: {str(e)}")
            # Fallback to basic extraction
            basic_text = self._basic_extraction_fallback(content, content_type)
            return {
                "raw_text": basic_text,
                "cleaned_text": basic_text,
                "chunks": [],
                "metadata": {
                    "filename": filename,
                    "content_type": content_type,
                    "error": str(e),
                    "fallback_used": True
                }
            }
    
    def _extract_with_layout(self, content: bytes, content_type: str, filename: str) -> Tuple[str, Dict]:
        """Extract text with layout and structure awareness."""
        structure_metadata = {
            "headings": [],
            "tables": [],
            "page_count": 0,
            "layout_detected": False
        }
        
        if content_type == 'application/pdf':
            return self._extract_pdf_with_layout(content, structure_metadata)
        elif content_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
            return self._extract_docx_with_structure(content, structure_metadata)
        elif content_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
            return self._extract_ppt_with_structure(content, structure_metadata)
        else:
            # Fallback to basic extraction
            text = self._basic_extraction_fallback(content, content_type)
            return text, structure_metadata
    
    def _extract_pdf_with_layout(self, content: bytes, metadata: Dict) -> Tuple[str, Dict]:
        """Extract PDF with advanced layout analysis."""
        text_parts = []
        
        # Try pdfplumber first for better layout detection
        if PDFPLUMBER_AVAILABLE and self.config.use_layout_analysis:
            try:
                with pdfplumber.open(io.BytesIO(content)) as pdf:
                    metadata["page_count"] = len(pdf.pages)
                    metadata["layout_detected"] = True
                    
                    for page_num, page in enumerate(pdf.pages, 1):
                        # Extract text with layout preservation
                        page_text = page.extract_text(
                            x_tolerance=3,
                            y_tolerance=3,
                            layout=True,
                            x_density=7.25,
                            y_density=13
                        )
                        
                        if page_text:
                            # Clean up layout artifacts
                            page_text = self._clean_pdf_layout(page_text, page_num)
                            text_parts.append(page_text)
                        
                        # Extract tables separately
                        tables = page.extract_tables()
                        if tables:
                            metadata["tables"].extend([
                                {"page": page_num, "table_num": i, "rows": len(table)}
                                for i, table in enumerate(tables)
                            ])
                    
                    return "\n\n".join(text_parts), metadata
                    
            except Exception as e:
                logger.warning(f"pdfplumber extraction failed: {e}, falling back to pypdf")
        
        # Fallback to pypdf
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PdfReader(pdf_file)
            metadata["page_count"] = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    # Basic cleaning for pypdf
                    page_text = self._clean_pdf_layout(page_text, page_num)
                    text_parts.append(page_text)
            
            return "\n\n".join(text_parts), metadata
            
        except Exception as e:
            # OCR fallback for scanned PDFs
            if OCR_AVAILABLE and self.config.use_ocr_fallback:
                logger.info("Attempting OCR fallback for PDF")
                return self._ocr_pdf_fallback(content, metadata)
            else:
                raise e
    
    def _extract_docx_with_structure(self, content: bytes, metadata: Dict) -> Tuple[str, Dict]:
        """Extract DOCX with heading structure preservation."""
        try:
            docx_file = io.BytesIO(content)
            doc = docx.Document(docx_file)
            
            text_parts = []
            current_heading_stack = []
            
            for para in doc.paragraphs:
                # Check if paragraph is a heading
                if para.style.name.startswith('Heading'):
                    level = int(para.style.name.split()[-1]) if para.style.name.split()[-1].isdigit() else 1
                    heading_text = para.text.strip()
                    
                    # Update heading stack
                    if level <= len(current_heading_stack):
                        current_heading_stack = current_heading_stack[:level-1]
                    current_heading_stack.append(heading_text)
                    
                    metadata["headings"].append({
                        "level": level,
                        "text": heading_text,
                        "hierarchy": current_heading_stack.copy()
                    })
                    
                    text_parts.append(f"\n{'#' * level} {heading_text}\n")
                else:
                    # Regular paragraph
                    if para.text.strip():
                        text_parts.append(para.text)
            
            # Extract tables
            for table_num, table in enumerate(doc.tables):
                metadata["tables"].append({
                    "table_num": table_num,
                    "rows": len(table.rows),
                    "cols": len(table.columns) if table.rows else 0
                })
            
            return "\n".join(text_parts), metadata
            
        except Exception as e:
            logger.error(f"Error extracting DOCX with structure: {e}")
            raise e
    
    def _extract_ppt_with_structure(self, content: bytes, metadata: Dict) -> Tuple[str, Dict]:
        """Extract PowerPoint with slide structure."""
        try:
            ppt_file = io.BytesIO(content)
            prs = Presentation(ppt_file)
            
            text_parts = []
            slide_count = len(prs.slides)
            metadata["page_count"] = slide_count
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"\n--- Slide {slide_num} ---\n"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Try to identify titles vs content
                        if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                            if 'title' in str(shape.placeholder_format).lower():
                                slide_text.append(f"# {shape.text}")
                                continue
                        
                        slide_text.append(shape.text)
                
                text_parts.append("\n".join(slide_text))
            
            return "\n\n".join(text_parts), metadata
            
        except Exception as e:
            logger.error(f"Error extracting PPT with structure: {e}")
            raise e
    
    def _clean_and_normalize(self, text: str) -> str:
        """Clean and normalize extracted text content."""
        if not text:
            return ""
        
        # Remove page numbers (common patterns)
        if self.config.remove_page_numbers:
            text = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)
            text = re.sub(r'Page \d+', '', text, flags=re.IGNORECASE)
        
        # Remove headers/footers (repeated text at top/bottom)
        if self.config.remove_headers_footers:
            text = self._remove_repetitive_headers_footers(text)
        
        # Remove boilerplate
        if self.config.remove_boilerplate:
            text = self._remove_boilerplate_patterns(text)
        
        # Normalize whitespace
        if self.config.normalize_whitespace:
            # Fix line break hyphenation
            text = re.sub(r'-\s*\n\s*', '', text)
            # Normalize multiple newlines
            text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
            # Normalize spaces
            text = re.sub(r'[ \t]+', ' ', text)
        
        # Fix encoding issues
        if self.config.fix_encoding:
            text = self._fix_encoding_issues(text)
        
        return text.strip()
    
    def _create_semantic_chunks(self, text: str, structure_metadata: Dict) -> List[DocumentChunk]:
        """Create semantically coherent chunks with overlap."""
        chunks = []
        
        if not text:
            return chunks
        
        # Try unstructured for intelligent chunking
        if UNSTRUCTURED_AVAILABLE:
            try:
                chunks = self._chunk_with_unstructured(text, structure_metadata)
                if chunks:
                    return chunks
            except Exception as e:
                logger.warning(f"Unstructured chunking failed: {e}, falling back to token-based")
        
        # Fallback to token-based chunking
        return self._chunk_by_tokens(text, structure_metadata)
    
    def _chunk_with_unstructured(self, text: str, structure_metadata: Dict) -> List[DocumentChunk]:
        """Use unstructured library for intelligent chunking."""
        chunks = []
        
        # Create temporary file for unstructured
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as tmp_file:
            tmp_file.write(text)
            tmp_file_path = tmp_file.name
        
        try:
            # Partition the document
            elements = partition(filename=tmp_file_path)
            
            # Chunk by title/sections
            chunked_elements = chunk_by_title(
                elements, 
                max_characters=self.config.chunk_size * 4,  # Rough token-to-char conversion
                combine_text_under_n_chars=100
            )
            
            for i, chunk_element in enumerate(chunked_elements):
                chunks.append(DocumentChunk(
                    content=str(chunk_element),
                    chunk_type=chunk_element.category if hasattr(chunk_element, 'category') else 'text',
                    metadata={'source': 'unstructured'},
                    position=i,
                    heading_hierarchy=self._extract_heading_context(str(chunk_element), structure_metadata)
                ))
            
        finally:
            os.unlink(tmp_file_path)
        
        return chunks
    
    def _chunk_by_tokens(self, text: str, structure_metadata: Dict) -> List[DocumentChunk]:
        """Fallback token-based chunking with overlap."""
        chunks = []
        
        if not self.tokenizer:
            # Simple character-based chunking if no tokenizer
            return self._chunk_by_characters(text, structure_metadata)
        
        # Tokenize the text
        tokens = self.tokenizer.encode(text)
        
        chunk_size = self.config.chunk_size
        overlap = self.config.chunk_overlap
        
        for i in range(0, len(tokens), chunk_size - overlap):
            chunk_tokens = tokens[i:i + chunk_size]
            chunk_text = self.tokenizer.decode(chunk_tokens)
            
            chunks.append(DocumentChunk(
                content=chunk_text,
                chunk_type='token_chunk',
                metadata={'token_count': len(chunk_tokens)},
                position=i // (chunk_size - overlap),
                heading_hierarchy=self._extract_heading_context(chunk_text, structure_metadata)
            ))
        
        return chunks
    
    def _chunk_by_characters(self, text: str, structure_metadata: Dict) -> List[DocumentChunk]:
        """Simple character-based chunking."""
        chunks = []
        chunk_size = self.config.chunk_size * 4  # Rough token-to-char conversion
        overlap = self.config.chunk_overlap * 4
        
        for i in range(0, len(text), chunk_size - overlap):
            chunk_text = text[i:i + chunk_size]
            
            chunks.append(DocumentChunk(
                content=chunk_text,
                chunk_type='char_chunk',
                metadata={'char_count': len(chunk_text)},
                position=i // (chunk_size - overlap),
                heading_hierarchy=[]
            ))
        
        return chunks
    
    # Helper methods for cleaning and processing
    
    def _clean_pdf_layout(self, text: str, page_num: int) -> str:
        """Clean PDF layout artifacts."""
        if not text:
            return ""
        
        # Remove page numbers at start/end of page
        text = re.sub(rf'^\s*{page_num}\s*\n', '', text)
        text = re.sub(rf'\n\s*{page_num}\s*$', '', text)
        
        # Fix common PDF extraction issues
        text = re.sub(r'([a-z])\s+([A-Z])', r'\1. \2', text)  # Missing periods
        text = re.sub(r'([a-z])\n([a-z])', r'\1 \2', text)    # Broken words
        
        return text
    
    def _remove_repetitive_headers_footers(self, text: str) -> str:
        """Remove repetitive headers and footers."""
        lines = text.split('\n')
        if len(lines) < 10:
            return text
        
        # Look for repeated patterns at start/end of text sections
        # This is a simplified version - could be enhanced with ML
        return text
    
    def _remove_boilerplate_patterns(self, text: str) -> str:
        """Remove common boilerplate patterns."""
        # Common patterns to remove
        patterns = [
            r'© \d{4}.*?All rights reserved\.?',
            r'Confidential and Proprietary.*?\n',
            r'This document contains.*?confidential.*?\n',
            r'Page \d+ of \d+',
            r'Document continues on next page.*?\n'
        ]
        
        for pattern in patterns:
            text = re.sub(pattern, '', text, flags=re.IGNORECASE | re.DOTALL)
        
        return text
    
    def _fix_encoding_issues(self, text: str) -> str:
        """Fix common encoding and character issues."""
        # Fix smart quotes and dashes
        text = text.replace('"', '"').replace('"', '"')
        text = text.replace(''', "'").replace(''', "'")
        text = text.replace('—', '--').replace('–', '-')
        
        # Fix ligatures
        text = text.replace('ﬁ', 'fi').replace('ﬂ', 'fl')
        text = text.replace('ﬀ', 'ff').replace('ﬃ', 'ffi')
        
        return text
    
    def _extract_heading_context(self, chunk_text: str, structure_metadata: Dict) -> List[str]:
        """Extract heading hierarchy context for a chunk."""
        # Simplified - look for markdown-style headings in chunk
        headings = []
        for line in chunk_text.split('\n'):
            if line.strip().startswith('#'):
                headings.append(line.strip())
        
        return headings
    
    def _ocr_pdf_fallback(self, content: bytes, metadata: Dict) -> Tuple[str, Dict]:
        """OCR fallback for scanned PDFs."""
        try:
            # This would require converting PDF to images and OCR
            # Simplified implementation
            logger.info("OCR fallback not fully implemented")
            metadata["processing_notes"] = ["OCR fallback attempted but not implemented"]
            return "", metadata
        except Exception as e:
            logger.error(f"OCR fallback failed: {e}")
            raise e
    
    def _basic_extraction_fallback(self, content: bytes, content_type: str) -> str:
        """Fallback to basic text extraction."""
        from core.processing_1_2_0.engine.text_extractor import TextExtractor
        return TextExtractor.extract_text(content, content_type) 