"""
Text extraction utilities for different file types.
"""

import io
from typing import Dict, Any
from bs4 import BeautifulSoup
import docx
from PyPDF2 import PdfReader
import pandas as pd
from pptx import Presentation
from openpyxl import load_workbook

class TextExtractor:
    """Handles text extraction from different file types."""
    
    @staticmethod
    def extract_text(content: bytes, content_type: str) -> str:
        """Extract text from content based on its type.
        
        Args:
            content: Raw content bytes
            content_type: MIME type of the content
            
        Returns:
            Extracted text content
        """
        content_type = content_type.lower()
        
        try:
            if content_type == 'application/pdf':
                return TextExtractor._extract_pdf_text(content)
            elif content_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
                return TextExtractor._extract_docx_text(content)
            elif content_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
                return TextExtractor._extract_ppt_text(content)
            elif content_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel', 'text/csv']:
                return TextExtractor._extract_excel_text(content, content_type)
            elif content_type == 'text/plain':
                return TextExtractor._extract_plain_text(content)
            elif content_type == 'text/html':
                return TextExtractor._extract_html_text(content)
            else:
                return f"Unsupported content type: {content_type}"
        except Exception as e:
            return f"Error extracting text: {str(e)}"
    
    @staticmethod
    def _extract_pdf_text(content: bytes) -> str:
        """Extract text from PDF content."""
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            return f"Error extracting PDF text: {str(e)}"
    
    @staticmethod
    def _extract_docx_text(content: bytes) -> str:
        """Extract text from DOCX content."""
        try:
            docx_file = io.BytesIO(content)
            doc = docx.Document(docx_file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            return f"Error extracting DOCX text: {str(e)}"
    
    @staticmethod
    def _extract_ppt_text(content: bytes) -> str:
        """Extract text from PowerPoint content."""
        try:
            ppt_file = io.BytesIO(content)
            prs = Presentation(ppt_file)
            text = []
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                if slide_text:
                    text.append("\n".join(slide_text))
            
            return "\n\n".join(text)
        except Exception as e:
            return f"Error extracting PowerPoint text: {str(e)}"
    
    @staticmethod
    def _extract_excel_text(content: bytes, content_type: str) -> str:
        """Extract text from Excel/CSV content."""
        try:
            if content_type == 'text/csv':
                # Handle CSV files
                csv_file = io.StringIO(content.decode('utf-8', errors='replace'))
                df = pd.read_csv(csv_file)
            else:
                # Handle Excel files
                excel_file = io.BytesIO(content)
                df = pd.read_excel(excel_file)
            
            # Convert DataFrame to string representation
            return df.to_string(index=False)
        except Exception as e:
            return f"Error extracting Excel/CSV text: {str(e)}"
    
    @staticmethod
    def _extract_plain_text(content: bytes) -> str:
        """Extract text from plain text content."""
        try:
            if isinstance(content, bytes):
                return content.decode('utf-8', errors='replace')
            return str(content)
        except Exception as e:
            return f"Error extracting plain text: {str(e)}"
    
    @staticmethod
    def _extract_html_text(content: bytes) -> str:
        """Extract text from HTML content."""
        try:
            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='replace')
            soup = BeautifulSoup(content, 'html.parser')
            return soup.get_text(separator='\n', strip=True)
        except Exception as e:
            return f"Error extracting HTML text: {str(e)}" 