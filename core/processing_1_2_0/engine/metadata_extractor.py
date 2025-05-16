"""
Metadata extraction utilities for different file types.
"""

import io
from typing import Dict, Any
import pandas as pd
from pptx import Presentation
from openpyxl import load_workbook

class MetadataExtractor:
    """Handles metadata extraction from different file types."""
    
    @staticmethod
    def extract_metadata(content: bytes, content_type: str) -> Dict[str, Any]:
        """Extract metadata from content based on its type.
        
        Args:
            content: Raw content bytes
            content_type: MIME type of the content
            
        Returns:
            Dictionary containing extracted metadata
        """
        content_type = content_type.lower()
        
        try:
            if content_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
                return MetadataExtractor._extract_ppt_metadata(content)
            elif content_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel', 'text/csv']:
                return MetadataExtractor._extract_excel_metadata(content, content_type)
            else:
                return {}
        except Exception as e:
            return {}
    
    @staticmethod
    def _extract_ppt_metadata(content: bytes) -> Dict[str, Any]:
        """Extract metadata from PowerPoint content."""
        try:
            ppt_file = io.BytesIO(content)
            prs = Presentation(ppt_file)
            
            metadata = {
                'title': prs.core_properties.title or '',
                'author': prs.core_properties.author or '',
                'last_modified': prs.core_properties.modified.isoformat() if prs.core_properties.modified else '',
                'slide_count': len(prs.slides),
                'slide_titles': []
            }
            
            # Extract slide titles
            for slide in prs.slides:
                title = ''
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        title = shape.text.strip()
                        break
                metadata['slide_titles'].append(title)
            
            return metadata
        except Exception as e:
            return {
                'title': '',
                'author': '',
                'last_modified': '',
                'slide_count': 0,
                'slide_titles': []
            }
    
    @staticmethod
    def _extract_excel_metadata(content: bytes, content_type: str) -> Dict[str, Any]:
        """Extract metadata from Excel/CSV content."""
        try:
            if content_type == 'text/csv':
                # Handle CSV files
                csv_file = io.StringIO(content.decode('utf-8', errors='replace'))
                df = pd.read_csv(csv_file)
                metadata = {
                    'sheet_count': 1,
                    'sheet_names': ['Sheet1'],
                    'row_count': len(df),
                    'column_count': len(df.columns),
                    'column_names': df.columns.tolist()
                }
            else:
                # Handle Excel files
                excel_file = io.BytesIO(content)
                wb = load_workbook(excel_file, read_only=True, data_only=True)
                metadata = {
                    'sheet_count': len(wb.sheetnames),
                    'sheet_names': wb.sheetnames,
                    'row_count': 0,
                    'column_count': 0,
                    'column_names': []
                }
                
                # Get dimensions of the first sheet
                if wb.sheetnames:
                    ws = wb[wb.sheetnames[0]]
                    metadata['row_count'] = ws.max_row
                    metadata['column_count'] = ws.max_column
                    metadata['column_names'] = [cell.value for cell in ws[1]] if ws.max_row > 0 else []
            
            return metadata
        except Exception as e:
            return {
                'sheet_count': 0,
                'sheet_names': [],
                'row_count': 0,
                'column_count': 0,
                'column_names': []
            } 