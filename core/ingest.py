import os
import sys
import asyncio
from dotenv import load_dotenv
from core.openai_service import openai_service
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd

# Load environment variables from .env
load_dotenv()

# Accept a file or directory as input (argument 1)
if len(sys.argv) < 2:
    print("Usage: python ingest.py <file-or-directory>")
    sys.exit(1)

input_path = sys.argv[1]

# Placeholder for upload logic
print(f"[INFO] Would ingest: {input_path}")

async def upload_and_print_file_id(file_path, metadata):
    file_id = await openai_service.upload_file_to_file_search(file_path, metadata)
    print(f"[UPLOAD] Uploaded {file_path} to File Search. file_id: {file_id}")

def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".txt":
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        elif ext == ".pdf":
            reader = PdfReader(file_path)
            return "\n".join(page.extract_text() or '' for page in reader.pages)
        elif ext == ".docx":
            doc = Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".pptx":
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        elif ext == ".csv":
            df = pd.read_csv(file_path, nrows=5)
            return f"Columns: {', '.join(df.columns)}\nSample rows:\n{df.to_string(index=False)}"
        elif ext in [".xlsx", ".xls"]:
            df = pd.read_excel(file_path, nrows=5)
            return f"Columns: {', '.join(df.columns)}\nSample rows:\n{df.to_string(index=False)}"
        else:
            print(f"[SKIP] Unsupported file type: {file_path}")
            return None
    except Exception as e:
        print(f"[WARN] Could not extract text from {file_path}: {e}")
        return None

def process_file(file_path):
    content = extract_text_from_file(file_path)
    if content:
        # Prepare metadata
        metadata = {
            "filename": os.path.basename(file_path),
            "filetype": os.path.splitext(file_path)[1].lower(),
            # Add more metadata fields as needed
        }
        # Upload to OpenAI File Search
        asyncio.run(upload_and_print_file_id(file_path, metadata))
        print(f"[INFO] First 200 chars: {content[:200]}")
        async def summarize_and_print(text):
            summary = await openai_service.summarize_text(text)
            print(f"[SUMMARY] {summary}")
        asyncio.run(summarize_and_print(content))

if os.path.isfile(input_path):
    process_file(input_path)
else:
    if os.path.isdir(input_path):
        print(f"[INFO] Ingesting directory: {input_path}")
        for entry in os.listdir(input_path):
            file_path = os.path.join(input_path, entry)
            if os.path.isfile(file_path):
                print(f"\n[INFO] Processing file: {file_path}")
                process_file(file_path)
            else:
                print(f"[SKIP] Not a file: {file_path}")
    else:
        print("[WARN] Directory ingestion not implemented yet.")
# TODO: Add file reading, conversion, and upload logic using openai_service 